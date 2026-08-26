#!/usr/bin/env python3
"""One-page deck: Hyperloom src/ code volume.

Design follows Hyperloom-correctness-code-volume.pptx and
Hyperloom-token-efficiency-PR1233-EN.pptx: 13.33x7.5in dark canvas, Consolas for
figures and labels, Segoe UI for prose, amber for the number that carries the
point.

Numbers come from loc-census.json / loc-redundancy.json / loc-scenarios.json.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(0x0A, 0x0F, 0x18)
CARD = RGBColor(0x0E, 0x14, 0x1E)
CARD_LINE = RGBColor(0x1E, 0x28, 0x38)
RULE = RGBColor(0x2A, 0x33, 0x42)
TEXT = RGBColor(0xE9, 0xED, 0xF2)
BODY = RGBColor(0xC7, 0xD0, 0xDC)
MUTED = RGBColor(0x7A, 0x84, 0x97)
DIM = RGBColor(0x3B, 0x47, 0x60)
GOLD = RGBColor(0xD8, 0xA6, 0x6A)
BLUE = RGBColor(0x86, 0xB7, 0xD6)
BLUE_HDR = RGBColor(0x8F, 0xBF, 0xE0)
BLUE_DIM = RGBColor(0x5E, 0x7C, 0x93)
GREEN = RGBColor(0x8F, 0xE0, 0xAE)
MONO = "Consolas"
SANS = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def rect(x, y, w, h, fill, line=None, rounded=False, adj=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if rounded and adj is not None:
        shape.adjustments[0] = adj
    return shape


def card(x, y, w, h):
    return rect(x, y, w, h, CARD, CARD_LINE, rounded=True, adj=0.04)


def text(x, y, w, h, runs, size=8, align=PP_ALIGN.LEFT, spacing=None):
    """runs: list of (string, size, font, bold, color) or a list of such lists."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    paragraphs = runs if runs and isinstance(runs[0], list) else [runs]
    for i, para_runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        for s, sz, font, bold, color in para_runs:
            r = p.add_run()
            r.text = s
            r.font.size = Pt(sz)
            r.font.name = font
            r.font.bold = bold
            r.font.color.rgb = color
    return box


def bar(x, y, w_total, w_fill, h=0.17, fill=GOLD, back=None):
    if back is not None:
        rect(x, y, w_total, h, back, rounded=True, adj=0.3)
    if w_fill > 0:
        rect(x, y, max(w_fill, 0.02), h, fill, rounded=True, adj=0.3)


# ---------------------------------------------------------------- header
rect(0, 0, 13.333, 7.5, BG)
text(
    0.50,
    0.14,
    11.20,
    0.40,
    [
        ("Code Volume", 18, SANS, True, TEXT),
        ("   what the tree is, and what could actually go", 18, SANS, False, MUTED),
    ],
)
text(
    9.90,
    0.10,
    2.93,
    0.40,
    [("AMD-AGI/Hyperloom · src/hyperloom @ 66e7c631", 7, MONO, False, DIM)],
    align=PP_ALIGN.RIGHT,
)
rect(0.50, 0.52, 12.33, 0.02, RULE)
text(
    0.50,
    0.55,
    12.33,
    0.24,
    [
        ("1,104 Python files", 8, MONO, False, MUTED),
        ("   ·   ", 8, MONO, False, DIM),
        ("561,407 physical", 8, MONO, False, MUTED),
        ("   ·   ", 8, MONO, False, DIM),
        ("385,997 SLOC", 8, MONO, False, MUTED),
        ("   ·   ", 8, MONO, False, DIM),
        ("262,393 logical lines (radon LLOC)", 8, MONO, True, GOLD),
    ],
)

# ---------------------------------------------------------------- KPI cards
kpis = [
    (
        0.50,
        "THE TREE TODAY",
        "262,393",
        "  logical",
        GOLD,
        "118,143 production  ·  144,250 test  —  tests outweigh product code 1.22 : 1",
    ),
    (
        4.70,
        "NO FEATURE REMOVED",
        "99.0%",
        "  must stay",
        BLUE,
        "only 5,744 physical lines are provably redundant — 1.0% of the tree",
    ),
    (
        8.90,
        "CORE FEATURES ONLY",
        "92,031",
        "  logical",
        GOLD,
        "77.9% of production code is the core loop; 22.1% is optional",
    ),
]
for x, label, big, unit, color, sub in kpis:
    card(x, 0.86, 3.92, 1.44)
    text(x + 0.20, 0.99, 3.52, 0.22, [(label, 7.8, MONO, True, MUTED)])
    text(
        x + 0.20,
        1.24,
        3.52,
        0.62,
        [(big, 34, MONO, True, color), (unit, 11, MONO, True, MUTED)],
    )
    text(x + 0.20, 1.92, 3.52, 0.30, [(sub, 8, SANS, False, BODY)])

# ---------------------------------------------------------------- panel: three ways to count
card(0.50, 2.48, 6.05, 1.86)
text(
    0.70,
    2.59,
    5.65,
    0.22,
    [
        ("THREE WAYS TO COUNT", 7.8, MONO, True, BLUE_HDR),
        ("    same tree, three defensible totals", 7.8, MONO, False, MUTED),
    ],
)
rows = [
    ("physical lines", "561,407", "278,845", "282,562", MUTED),
    ("SLOC  no blank / comment / docstring", "385,997", "176,423", "209,574", MUTED),
    ("logical lines  radon LLOC", "262,393", "118,143", "144,250", GOLD),
    ("AST statements", "218,595", "97,049", "121,546", MUTED),
]
text(0.70, 2.86, 2.60, 0.20, [("definition", 6.6, MONO, False, DIM)])
for j, head in enumerate(("all", "prod", "test")):
    text(
        3.55 + j * 0.95,
        2.86,
        0.90,
        0.20,
        [(head, 6.6, MONO, False, DIM)],
        align=PP_ALIGN.RIGHT,
    )
for i, (name, a, p, t, color) in enumerate(rows):
    y = 3.06 + i * 0.24
    bold = color is GOLD
    text(0.70, y, 2.90, 0.20, [(name, 7.2, MONO, bold, BODY if not bold else GOLD)])
    for j, v in enumerate((a, p, t)):
        text(
            3.55 + j * 0.95,
            y,
            0.90,
            0.20,
            [(v, 7.4, MONO, True, color)],
            align=PP_ALIGN.RIGHT,
        )
text(
    0.70,
    4.06,
    5.65,
    0.20,
    [
        ("31.2%", 7.4, SANS, True, GOLD),
        (" of physical lines are blank, comment or docstring.", 7.2, SANS, False, BODY),
    ],
)

# ---------------------------------------------------------------- panel: where the mass is
card(6.78, 2.48, 6.05, 1.86)
text(
    6.98,
    2.59,
    3.60,
    0.22,
    [
        ("WHERE THE MASS IS", 7.8, MONO, True, BLUE_HDR),
        ("    production only", 7.8, MONO, False, MUTED),
    ],
)
text(
    10.20,
    2.59,
    2.63,
    0.22,
    [
        ("core", 7, MONO, False, MUTED),
        ("  ■", 7, MONO, True, BLUE_DIM),
        ("      optional", 7, MONO, False, MUTED),
        ("  ■", 7, MONO, True, GOLD),
    ],
    align=PP_ALIGN.RIGHT,
)
features = [
    ("agents/kernel", 18183, True),
    ("orchestrator/actions", 16728, True),
    ("io/breakdown  report", 11337, True),
    ("orchestrator/phases", 10122, True),
    ("orchestrator/kernel", 7438, True),
    ("orchestrator/loop", 6925, True),
    ("agents/robustness", 5209, False),
    ("io/multi_node", 5011, False),
]
SCALE = 2.35 / 18183
for i, (name, lloc, is_core) in enumerate(features):
    y = 2.88 + i * 0.155
    text(6.98, y - 0.01, 1.85, 0.18, [(name, 6.6, MONO, False, BODY if is_core else GOLD)])
    bar(8.90, y, 2.35, lloc * SCALE, h=0.12, fill=BLUE_DIM if is_core else GOLD)
    text(
        11.35,
        y - 0.02,
        0.72,
        0.18,
        [(f"{lloc:,}", 6.6, MONO, True, BLUE if is_core else GOLD)],
        align=PP_ALIGN.RIGHT,
    )
text(
    6.98,
    4.14,
    5.65,
    0.20,
    [
        ("27 further features", 7, SANS, True, TEXT),
        ("  hold 26,870 core and 10,320 optional logical lines.", 7, SANS, False, BODY),
    ],
)

# ---------------------------------------------------------------- panel: provably redundant
card(0.50, 4.42, 6.05, 2.06)
text(
    0.70,
    4.53,
    5.65,
    0.22,
    [
        ("KEEP EVERY FEATURE", 7.8, MONO, True, BLUE_HDR),
        ("    what is provably safe to delete", 7.8, MONO, False, MUTED),
    ],
)
red_rows = [
    ("token-identical files", "16 files", "807", GREEN),
    ("token-identical function bodies", "475 copies", "4,937", GREEN),
    ("provably removable", "1.0% of tree", "5,744", GOLD),
]
for i, (name, mid, val, color) in enumerate(red_rows):
    y = 4.82 + i * 0.24
    bold = color is GOLD
    if bold:
        rect(0.70, y - 0.03, 5.65, 0.01, RULE)
        y += 0.02
    text(0.70, y, 2.60, 0.20, [(name, 7.2, MONO, bold, GOLD if bold else BODY)])
    text(3.35, y, 1.55, 0.20, [(mid, 7, MONO, False, MUTED)])
    text(
        5.05,
        y,
        1.30,
        0.20,
        [(val, 7.6, MONO, True, color), (" lines", 6.6, MONO, False, MUTED)],
        align=PP_ALIGN.RIGHT,
    )
text(
    0.70,
    5.62,
    5.65,
    0.62,
    [
        [
            ("Dead code does not survive review — ", 7, SANS, True, GOLD),
            (
                "vulture flags 4,938 lines, but all 15 of the largest are reached through a string dispatch table or a self-registering renderer registry. Real dead code is under a hundred lines.",
                7,
                SANS,
                False,
                BODY,
            ),
        ],
        [
            ("Coverage-chasing tests are the one large lever — ", 7, SANS, True, GOLD),
            (
                "106 files named _units / _branches / _coverage hold 50,535 lines. Deleting them costs coverage, not a single product feature.",
                7,
                SANS,
                False,
                BODY,
            ),
        ],
    ],
    spacing=3,
)

# ---------------------------------------------------------------- panel: core only
card(6.78, 4.42, 6.05, 2.06)
text(
    6.98,
    4.53,
    5.65,
    0.22,
    [
        ("CORE FEATURES ONLY", 7.8, MONO, True, BLUE_HDR),
        ("    TraceLens · Arbor loop · GEAK · report · CLI", 7.8, MONO, False, MUTED),
    ],
)
core_rows = [
    ("production code today", "487 files", "118,143", MUTED),
    ("drop 6 optional features", "−152 files", "−20,540", GOLD),
    ("trim optional surface inside core", "−24 files", "−5,572", GOLD),
    ("core needs", "311 files", "92,031", GREEN),
]
for i, (name, mid, val, color) in enumerate(core_rows):
    y = 4.82 + i * 0.24
    last = i == len(core_rows) - 1
    if last:
        rect(6.98, y - 0.03, 5.65, 0.01, RULE)
        y += 0.02
    text(6.98, y, 3.05, 0.20, [(name, 7.2, MONO, last, GREEN if last else BODY)])
    text(10.10, y, 1.20, 0.20, [(mid, 7, MONO, False, MUTED)], align=PP_ALIGN.RIGHT)
    text(
        11.40,
        y,
        1.23,
        0.20,
        [(val, 7.6, MONO, True, color), (" lloc", 6.6, MONO, False, MUTED)],
        align=PP_ALIGN.RIGHT,
    )
text(
    6.98,
    5.78,
    5.65,
    0.46,
    [
        [
            ("Dropped: ", 7, SANS, True, GOLD),
            (
                "robustness 5,209 · multi_node 5,011 · framework agent 5,698 · critic 2,260 · quantization 782 · Forge backend 2,825 · collective codegen 1,122 · diffusion 684 · FlyDSL 300 · 13 report sections 641",
                7,
                SANS,
                False,
                BODY,
            ),
        ],
        [
            ("The cost: ", 7, SANS, True, GOLD),
            (
                "21 feature-pair import edges cross the line. The framework agent is the only one needing real refactoring — four core features import it.",
                7,
                SANS,
                False,
                BODY,
            ),
        ],
    ],
    spacing=3,
)

# ---------------------------------------------------------------- conclusion
card(0.50, 6.56, 12.33, 0.50)
text(
    0.70,
    6.68,
    12.00,
    0.32,
    [
        ("Conclusion    ", 8, MONO, True, BLUE),
        ("This tree is not mostly optional features — ", 8.6, SANS, True, GREEN),
        (
            "cutting every non-core feature removes 22%; keeping them all allows 1%. The volume sits in the test tree and in lines that are not code.",
            8.6,
            SANS,
            True,
            BODY,
        ),
    ],
)
text(
    0.50,
    7.20,
    12.33,
    0.18,
    [
        (
            "Method: radon raw LLOC/SLOC · AST statements · token-normalised sha256 duplicate detection (file + function, overlap removed) · vulture 60% reviewed item by item · import graph from 53 entry points.  Reproduce: loc-census.py, loc-redundancy.py, loc-scenarios.py",
            6.5,
            MONO,
            False,
            DIM,
        )
    ],
)

out = "Hyperloom-loc-census-EN.pptx"
prs.save(out)
print("wrote", out)
