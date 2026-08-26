#!/usr/bin/env python3
"""Render a pptx slide to PNG with matplotlib, for layout checking only.

Approximates text metrics; use it to catch overlap and overflow, not to judge
typography.
"""
from __future__ import annotations

import sys

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.util import Emu

src = sys.argv[1] if len(sys.argv) > 1 else "Hyperloom-loc-census-EN.pptx"
out = sys.argv[2] if len(sys.argv) > 2 else "loc-slide-preview.png"

prs = Presentation(src)
W = Emu(prs.slide_width).inches
H = Emu(prs.slide_height).inches
slide = prs.slides[0]

fig, ax = plt.subplots(figsize=(W, H), dpi=110)
ax.set_xlim(0, W)
ax.set_ylim(H, 0)
ax.axis("off")
fig.subplots_adjust(0, 0, 1, 1)


def rgb(c):
    return f"#{c}" if isinstance(c, str) else f"#{c!s}"


for sh in slide.shapes:
    x, y = Emu(sh.left).inches, Emu(sh.top).inches
    w, h = Emu(sh.width).inches, Emu(sh.height).inches
    if sh.shape_type == 1:  # AUTO_SHAPE
        try:
            fc = "#" + str(sh.fill.fore_color.rgb)
        except Exception:
            fc = "none"
        try:
            ec = "#" + str(sh.line.color.rgb)
        except Exception:
            ec = "none"
        rounded = "Rounded" in sh.name
        patch_cls = FancyBboxPatch if rounded else Rectangle
        if rounded:
            p = FancyBboxPatch(
                (x + 0.02, y + 0.02),
                w - 0.04,
                h - 0.04,
                boxstyle="round,pad=0.02,rounding_size=0.05",
                facecolor=fc,
                edgecolor=ec if ec != "none" else "none",
                linewidth=0.6,
            )
        else:
            p = Rectangle((x, y), w, h, facecolor=fc, edgecolor="none")
        ax.add_patch(p)
        continue

    if not sh.has_text_frame:
        continue
    # Draw each paragraph as one line of concatenated runs.
    cy = y
    for para in sh.text_frame.paragraphs:
        runs = [r for r in para.runs if r.text]
        if not runs:
            cy += 0.12
            continue
        sizes = [r.font.size.pt if r.font.size else 8 for r in runs]
        line_h = max(sizes) / 72 * 1.35
        cx = x
        align = str(para.alignment)
        full = "".join(r.text for r in runs)
        if "RIGHT" in align:
            est = sum(len(r.text) * (s * 0.55) / 72 for r, s in zip(runs, sizes))
            cx = x + w - est
        for r, s in zip(runs, sizes):
            try:
                col = "#" + str(r.font.color.rgb)
            except Exception:
                col = "#C7D0DC"
            mono = (r.font.name or "").lower().startswith("cons")
            ax.text(
                cx,
                cy + line_h * 0.78,
                r.text,
                fontsize=s * 0.98,
                color=col,
                family="monospace" if mono else "sans-serif",
                weight="bold" if r.font.bold else "normal",
                va="baseline",
                ha="left",
            )
            cx += len(r.text) * (s * (0.55 if mono else 0.5)) / 72
        # Wrap indicator: mark text that runs past its own box.
        if cx > x + w + 0.05:
            ax.plot([x + w, x + w], [cy, cy + line_h], color="#FF3B30", lw=1.2)
        cy += line_h + 0.02
        _ = full

fig.savefig(out, facecolor="#0A0F18")
print("wrote", out)
